import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

def main():
    prs = Presentation()
    
    # 1. 設定投影片尺寸 13.33" × 7.50" (LAYOUT_WIDE)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.50)
    
    # 2. 定義顏色常數
    COLOR_DARK_COVER = RGBColor(0x0A, 0x16, 0x28) # #0A1628 (Cover 暗底)
    COLOR_TEAL_FILL = RGBColor(0x12, 0x5C, 0x69)  # #125C69 (38x fills 深松石綠)
    COLOR_LIGHT_CARD = RGBColor(0xDB, 0xEF, 0xEF) # #DBEFEF (10x cards 淺松石綠)
    COLOR_TEXT = RGBColor(0x33, 0x41, 0x55)       # #334155 (主要文字深灰藍)
    COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)      # 白色背景
    
    # 3. 取得空白版面 (Layout 6 通常是無預設佔位符的空白版型)
    blank_layout = prs.slide_layouts[6]
    
    # ==================== SLIDE 1: 封面 (暗底) ====================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # 設定暗底背景色
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_COVER
    
    # 封面主要文字框
    cover_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.5))
    tf1 = cover_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
    
    # 主標題
    p1 = tf1.paragraphs[0]
    p1.text = "Antigravity"
    p1.font.name = 'Abadi'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(20)
    
    # 副標題
    p2 = tf1.add_paragraph()
    p2.text = "引領 MAT 團隊的 Agentic AI 協作新時代"
    p2.font.name = '微軟正黑體'
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_LIGHT_CARD
    p2.space_after = Pt(48)
    
    # 署名
    p3 = tf1.add_paragraph()
    p3.text = "MAT 團隊分享會  |  2026.06.01"
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_LIGHT_CARD
    
    # ==================== 共用元素生成函數 ====================
    def add_slide_decorations(slide, section_text, title_text):
        # 設定白色底色背景 (防跑版)
        bg = slide.background
        bf = bg.fill
        bf.solid()
        bf.fore_color.rgb = COLOR_WHITE

        # 1. Accent Bar: 0.08" × 0.52" @ (0.7", 0.5")
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.5), Inches(0.08), Inches(0.52)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = COLOR_TEAL_FILL
        accent.line.fill.background() # 無邊框
        
        # 2. Title: 頂邊界 0.45", 左邊界 0.9" (位於 Accent Bar 右側)
        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.42), Inches(11.5), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = '微軟正黑體'
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEAL_FILL
        
        # 3. Section Pill: 3.70" × 0.42" @ (0.6", 1.3")
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.3), Inches(3.70), Inches(0.42)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = COLOR_LIGHT_CARD
        pill.line.fill.background()
        
        ptf = pill.text_frame
        ptf.word_wrap = True
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ptf.margin_left = Inches(0.18)
        ptf.margin_right = Inches(0.18)
        ptf.margin_top = ptf.margin_bottom = 0
        
        pp = ptf.paragraphs[0]
        pp.text = section_text
        pp.font.name = 'Arial'
        pp.font.size = Pt(12)
        pp.font.bold = True
        pp.font.color.rgb = COLOR_TEAL_FILL
        pp.alignment = PP_ALIGN.LEFT

    # ==================== 卡片加入函數 ====================
    def add_card(slide, left, top, width, height, title, content_lines):
        # 卡片底色 (圓角矩形)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT_CARD
        card.line.fill.background()
        
        # 文字方塊覆蓋在上面，提供完美的排版微調
        tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.25), width - Inches(0.5), height - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # 卡片標題
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = '微軟正黑體'
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEAL_FILL
        p_title.space_after = Pt(14)
        
        # 卡片內容
        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.name = '微軟正黑體'
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(8)
            
    # ==================== SLIDE 2: 什麼是 Antigravity？ ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide2, "01 / 核心概念", "什麼是 Antigravity？")
    
    add_card(
        slide2, 
        Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
        "Agentic AI 的自主開發體驗",
        [
            "• 自主規劃能力：動手開發前，自主理解並撰寫實作計劃，徹底告別傳統 AI 的盲目猜測。",
            "• 自我除錯與執行：在安全沙盒中執行 PowerShell 指令、自力編寫並測試代碼，直到正確為止。",
            "• 協同工作流：它不只是一個程式編輯器，更是一位能提供專業解決方案的「虛擬開發隊友」。"
        ]
    )
    
    add_card(
        slide2, 
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
        "極致踏實的「三大規劃機制」",
        [
            "• 實作計畫 (Implementation Plan)：將模糊的業務邏輯細化為具體技術步驟，取得你的同意後才動手。",
            "• 任務清單 (task.md)：透過即時 TODO 狀態條，讓開發過程中的每一步進度皆清晰、公開且透明。",
            "• 成果導覽 (walkthrough.md)：在開發完成後，自動演示測試過程與最終成效，提供視覺化的驗證回報。"
        ]
    )

    # ==================== SLIDE 3: 環境設定極速上手 ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide3, "02 / 環境設定", "MAT 同事的環境防污染指南")
    
    add_card(
        slide3,
        Inches(0.6), Inches(2.0), Inches(3.7), Inches(4.8),
        "01 / 建立獨立工作區",
        [
            "• 專屬沙盒：專案建議統一建在 scratch 下的合理子目錄。",
            "• 專案 Active Workspace：使用 VS Code 打開該子目錄，提供 AI 完整的上下文分析視野。"
        ]
    )
    
    add_card(
        slide3,
        Inches(4.7), Inches(2.0), Inches(3.7), Inches(4.8),
        "02 / 乾淨虛擬環境 (硬性)",
        [
            "• 堅持使用 uv 與 .venv：嚴格執行 Python 環境管理，絕不在 Base 環境下跑 Python。",
            "• 隔離風險：避免不同專案的套件與相依性衝突，確保您的日常分析工具穩定不損毀。"
        ]
    )
    
    add_card(
        slide3,
        Inches(8.8), Inches(2.0), Inches(3.7), Inches(4.8),
        "03 / 自動版控與分支安全",
        [
            "• gh CLI 自動化：首個版本會自動使用 GitHub CLI 建立遠端儲存庫並完成上傳。",
            "• 分支管理：若有重大變動，AI 會自主開 Branch 進行安全實作，完成後再詢問合併。"
        ]
    )

    # ==================== SLIDE 4: 自訂協作規範 ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide4, "02 / 自訂規範", "MAT 團隊的 AI 協作自訂規範 (Global Rules)")
    
    add_card(
        slide4,
        Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
        "MAT 團隊目前採用的全域偏好 (Rules)",
        [
            "• 繁中與台灣用語：討論過程、內部思考、程式註解、變更報告與 Commit Message 一律使用台灣繁體中文用語。",
            "• 隔離虛擬環境：開發 Python 程式時，強制使用 uv 搭配 .venv 管理環境，絕不在 Base 下跑 Python，保護主機環境。",
            "• 版控與分支安全：首版自動使用 gh 建立公開儲存庫上傳，有重大功能異動自動切換獨立分支，確認後再合併。"
        ]
    )
    
    add_card(
        slide4,
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
        "如何讓大家的 Antigravity 都套用此規範？",
        [
            "• 方法 A ── 專案級 GEMINI.md (最推薦！)：在專案根目錄下建立 GEMINI.md 檔案 (亦相容 CLAUDE.md)，寫入偏好命令。Antigravity 載入專案時會自動遵守，極利於團隊協作分享。",
            "• 方法 B ── 全域級 GEMINI.md (一勞永逸！)：在個人家目錄的隱藏設定資料夾下，建立 C:\\Users\\<Username>\\.gemini\\GEMINI.md 檔案。此後全電腦所有專案都會自動套用此規則！",
            "• 方法 C ── 系統 User Rules 設定：直接在 Antigravity 軟體設定的 User Rules 欄位中貼上這些 Rules，即可跨專案自動套用您的個人開發偏好。"
        ]
    )

    # ==================== SLIDE 5: 5分鐘現場 Demo 流程 ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide5, "03 / 實戰展示", "Demo 現場：5分鐘見證 AI 自主開發")
    
    # 4個步驟小卡片
    add_card(
        slide5,
        Inches(0.6), Inches(2.0), Inches(2.75), Inches(4.8),
        "1. 引導對齊",
        [
            "輸入 /brainstorm 指令啟動需求對齊。",
            "故意選擇「小白模式」，展示 AI 如何用大白話梳理需求、列出理解假設並提供方案比較表。"
        ]
    )
    
    add_card(
        slide5,
        Inches(3.6), Inches(2.0), Inches(2.75), Inches(4.8),
        "2. 核准實作計畫",
        [
            "確認 AI 產出的實作計畫書無誤後，輸入「同意」一鍵核准執行。",
            "AI 隨即生成 task.md 開始追蹤工作進度。"
        ]
    )
    
    add_card(
        slide5,
        Inches(6.6), Inches(2.0), Inches(2.75), Inches(4.8),
        "3. 自律虛擬執行",
        [
            "AI 自動在背景使用 uv 建立虛擬環境、下載相依套件，並自力編寫核心代碼與 Error Handling。"
        ]
    )
    
    add_card(
        slide5,
        Inches(9.6), Inches(2.0), Inches(2.75), Inches(4.8),
        "4. 成果驗證",
        [
            "AI 自動跑測試並撰寫 walkthrough.md，將成果與圖表以 Carousel 或 Markdown 完美打包呈獻給您。"
        ]
    )

    # ==================== SLIDE 6: 動手前先想清楚：/brainstorm ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide6, "04 / 亮點功能 01", "動手前先想清楚：/brainstorm")
    
    add_card(
        slide6,
        Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
        "🟢 小白模式 (推薦非技術同事)",
        [
            "• 成果導向：完全過濾艱澀的技術名詞。例如以「自動執行小程式」替代「Cron Job」。",
            "• 專注價值：對話重點聚焦於「你能得到什麼」、「有什麼使用風險」以及「未來可以如何延伸」。",
            "• 直觀對比：方案比較表簡化為優缺點、需要你配合做的事，讓決策輕鬆無負擔。"
        ]
    )
    
    add_card(
        slide6,
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8),
        "🔴 工程師模式 (高效率開發)",
        [
            "• 技術對齊：直接探討檔案路徑、套件版本、API 串接細節與伺服器架構，絕不拖泥帶水。",
            "• 極速開發：用最精簡的術語與 AI 同頻共振，快速切入代碼實作。",
            "• 技術矩陣：方案比較表列出完整的技術棧對比、運作效能與工時預估。"
        ]
    )

    # ==================== SLIDE 7: 打造 MAT 專屬的 AI 自動化工具庫 ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide7, "05 / 亮點功能 02", "workflow-skill-creator：讓經驗重複使用")
    
    # 左側大卡片
    add_card(
        slide7,
        Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.8),
        "做過一次，學會一輩子",
        [
            "• 經驗資產化：當你與 AI 合作成功完成了一套繁琐的實驗數據清洗、報告生成或文件處理流程，直接對它說「把我們剛剛做的事情建立成一個 Skill」。",
            "• 永久複用：AI 會將此流程提煉、打包成結構化的指令與指令腳本。下次只需一句話，AI 就能瞬間以完全相同的標準為您重複執行該任務。"
        ]
    )
    
    # 右上卡片
    add_card(
        slide7,
        Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.2),
        "工業級防護護欄 ── Rate Limiting",
        [
            "• 自動實作指數退避 (Exponential Backoff)，呼叫內部系統 API 時不怕造成負載，安全防封鎖。"
        ]
    )
    
    # 右下卡片
    add_card(
        slide7,
        Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.2),
        "跨進程鎖 ── File-lock 安全防衝突",
        [
            "• 多個 AI 協同運行時自動共用進程鎖，避免讀寫衝突與代碼覆寫，展現強健的開發品質。"
        ]
    )

    # ==================== SLIDE 8: 如何與 AI 隊友愉快協作？ ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide8, "06 / 協作指南", "Antigravity 的協作最佳實踐")
    
    add_card(
        slide8,
        Inches(0.6), Inches(2.0), Inches(3.7), Inches(4.8),
        "大招 /goal 指令",
        [
            "• 終極耐心：當您有極其複雜、需要反覆嘗試、或希望跑一整夜自動化測試的任務時，輸入 /goal 指令啟動。",
            "• 任務死守：AI 將進入徹底貫徹目標的深跑模式，不達目的絕不罷休。"
        ]
    )
    
    add_card(
        slide8,
        Inches(4.7), Inches(2.0), Inches(3.7), Inches(4.8),
        "神器 /browser 指令",
        [
            "• 線上檢索：當開發涉及最新推出的科學文獻、API 文件或遇到極難排除的 Bug 時，引導 AI 啟動瀏覽功能。",
            "• 掌握即時：實時上網翻閱官方文件，避免使用過時資訊開發代碼。"
        ]
    )
    
    add_card(
        slide8,
        Inches(8.8), Inches(2.0), Inches(3.7), Inches(4.8),
        "習慣 台灣繁體中文版控",
        [
            "• 道地用語：Antigravity 撰寫的程式註解、commit message、以及報告文件皆使用道地台灣用語，無縫融入團隊溝通。",
            "• 主分支防禦：開發過程嚴格於獨立分支進行，確保主程式庫的極致純淨。"
        ]
    )

    # 4. 存檔
    output_filename = "MAT_Antigravity_Introduction_v4.pptx"
    prs.save(output_filename)
    print(f"Success! PPTX file generated successfully: {output_filename}")

if __name__ == "__main__":
    main()
